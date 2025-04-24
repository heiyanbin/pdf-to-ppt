import os
import json
import click
from datetime import datetime
from typing import List
#from PyPDF2 import PdfReader
from llama_parse import LlamaParse
from pptx import Presentation
from dotenv import load_dotenv
from pydantic import BaseModel
from openai import OpenAI

class Slide(BaseModel):
    title: str
    content: List[str]
    layout_type: str = "title_and_content"

class SlidesResponse(BaseModel):
    slides: List[Slide]

load_dotenv()

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF file"""
    #with open(pdf_path, 'rb') as file:
    #    reader = PdfReader(file)
    #    text = "\n".join([page.extract_text() for page in reader.pages])
    
    documents = LlamaParse(result_type="text").load_data(pdf_path)
    combined_text = "\n\n".join(doc.text for doc in documents)
    # Save extracted text to content.txt
    try:
        with open('content.txt', 'w') as f:
            f.write(combined_text)
    except Exception as e:
        print(f"Warning: Could not save extracted text - {str(e)}")
    
    return combined_text

def generate_slide_structure(text):
    """Generate PPT slide structure using OpenAI"""
    prompt = f"""
    Extract key points from this content and create a ppt outline. the outline should be in json format.
    Content: {text}
    """
    client = OpenAI(base_url=os.getenv("OPENAI_API_BASE"))
    response = client.beta.chat.completions.parse(
        model=os.getenv("OPENAI_MODEL_NAME"),
        messages=[{"role": "user", "content": prompt}],
        temperature=0.6,  # Higher temperature for more creative/verbose output
        max_tokens=3000,  # Increased token limit for longer responses
        response_format=SlidesResponse
    )
    
    try:
        slides_data = json.loads(response.choices[0].message.content)
        # Save slides data to outline.json
        try:
            with open('outline.json', 'w') as f:
                json.dump(slides_data, f, indent=2)
        except Exception as e:
            print(f"Warning: Could not save outline - {str(e)}")
        return slides_data
    except json.JSONDecodeError:
        raise ValueError("Failed to parse OpenAI response as JSON")
    except Exception as e:
        raise ValueError(f"Response validation failed: {str(e)}")

def create_ppt(slides_data, output_path):
    """Generate PPT from slide structure"""
    prs = Presentation()
    
    for slide_data in slides_data['slides']:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = slide_data.get('title', 'Untitled')
        
        content = slide.placeholders[1]
        content.text = "\n".join(slide_data.get('content', []))
    
    prs.save(output_path)

@click.command()
@click.argument('input_file')
def main(input_file):
    """Convert PDF/TXT to PPT"""
    try:
        # Read input file
        if input_file.lower().endswith('.pdf'):
            text = extract_text_from_pdf(input_file)
        else:
            with open(input_file, 'r') as file:
                text = file.read()
        print('Document loaded successfully.') 
        # Generate slide structure
        slides = generate_slide_structure(text)
        print("Slide structure generated.")
        
        # Create PPT
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = f"output_{timestamp}.pptx"
        create_ppt(slides, output_path)
        
        print(f"PPT generated successfully: {output_path}")
    
    except Exception as e:
        print(f"Error: {str(e)}")

if __name__ == "__main__":
    main()
